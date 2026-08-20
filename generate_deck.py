import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette ──
NAVY = RGBColor(26, 39, 68)         # #1a2744
CRIMSON = RGBColor(192, 57, 43)     # #c0392b
MUTED = RGBColor(94, 110, 130)      # #5e6e82
LIGHT_BG = RGBColor(247, 248, 251)  # #f7f8fb
WHITE = RGBColor(255, 255, 255)
CARD_BORDER = RGBColor(219, 229, 232)
CARD_BG = RGBColor(255, 255, 255)
DARK_CARD = RGBColor(34, 49, 80)
ACCENT_BG = RGBColor(253, 240, 238)
GREEN = RGBColor(46, 125, 90)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, kicker_text, title_text, category=""):
    # Header container
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.1))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p_kicker = tf.paragraphs[0]
    p_kicker.text = kicker_text.upper()
    p_kicker.font.size = Pt(10)
    p_kicker.font.bold = True
    p_kicker.font.color.rgb = CRIMSON
    p_kicker.space_after = Pt(2)
    
    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = NAVY

def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

# ══════════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide (Dark Navy Background)
# ══════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
set_slide_background(s1, NAVY)

# Mascot image
if os.path.exists("public/mascot.png"):
    s1.shapes.add_picture("public/mascot.png", Inches(8.2), Inches(1.2), width=Inches(4.3))

tb1 = s1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(7.0), Inches(4.5))
tf1 = tb1.text_frame
tf1.word_wrap = True

p = tf1.paragraphs[0]
p.text = "SICKLESENSE"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.space_after = Pt(8)

p2 = tf1.add_paragraph()
p2.text = "Citation-Bound Clinical Evidence Assistant for Sickle Cell Disease"
p2.font.size = Pt(20)
p2.font.color.rgb = RGBColor(255, 138, 122)
p2.space_after = Pt(24)

p3 = tf1.add_paragraph()
p3.text = "A grounded RAG system integrating domain-specific BiomedBERT embeddings, Azure AI Search hybrid retrieval, and constrained LLM synthesis."
p3.font.size = Pt(13)
p3.font.color.rgb = RGBColor(203, 214, 226)
p3.space_after = Pt(36)

p4 = tf1.add_paragraph()
p4.text = "Presented by Team MINOR THREAT 2026 | Creativa SCD Hackathon"
p4.font.size = Pt(12)
p4.font.bold = True
p4.font.color.rgb = RGBColor(160, 180, 205)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 2: The Clinical Problem
# ══════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
set_slide_background(s2, LIGHT_BG)
add_header(s2, "The Clinical Challenge", "Why General-Purpose AI Fails in Sickle Cell Care")

# Card 1: Time Pressure & Information Overload
add_card(s2, 0.8, 1.6, 3.6, 5.2)
tb = s2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(3.2), Inches(4.8))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Fragmented Clinical Guidelines"
p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = NAVY; p.space_after = Pt(12)
p_body = tf.add_paragraph()
p_body.text = "• Authoritative protocols (NHLBI, WHO, regional studies) span hundreds of pages.\n\n• Clinicians and hematology trainees lack fast, trustworthy margin-note tools during urgent patient rounds.\n\n• Critical dosage and screening thresholds (e.g., TCD velocities, hydroxyurea titration) are buried in dense appendices."
p_body.font.size = Pt(12); p_body.font.color.rgb = MUTED

# Card 2: The Hallucination Danger
add_card(s2, 4.8, 1.6, 3.6, 5.2)
tb = s2.shapes.add_textbox(Inches(5.0), Inches(1.8), Inches(3.2), Inches(4.8))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Generic LLM Hallucinations"
p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = CRIMSON; p.space_after = Pt(12)
p_body = tf.add_paragraph()
p_body.text = "• Standard LLMs generate ungrounded, unverified medical advice with false confidence.\n\n• Fabricated citations and inaccurate contraindication advice create severe patient risk.\n\n• No page-level provenance or verification checks against the true medical corpus."
p_body.font.size = Pt(12); p_body.font.color.rgb = MUTED

# Card 3: Unsafe Personalized Prescribing
add_card(s2, 8.8, 1.6, 3.6, 5.2)
tb = s2.shapes.add_textbox(Inches(9.0), Inches(1.8), Inches(3.2), Inches(4.8))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Uncontrolled Prescribing Risk"
p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = NAVY; p.space_after = Pt(12)
p_body = tf.add_paragraph()
p_body.text = "• Population guidelines must never be confused with direct individual patient prescriptions.\n\n• Existing chatbots fail to intercept patient-specific inquiries (e.g. 'What dose should I give my child?').\n\n• Lack of strict safety interceptors creates clinical liability."
p_body.font.size = Pt(12); p_body.font.color.rgb = MUTED

# ══════════════════════════════════════════════════════════════════════
# SLIDE 3: The Solution - SickleSense
# ══════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
set_slide_background(s3, LIGHT_BG)
add_header(s3, "Our Solution", "SickleSense: The Citation-Bound Clinical Margin Note")

# 4 Pillars Grid
pillars = [
    ("Strict Grounding", "Every generated claim is bound to retrieved evidence chunks. Zero unverified external speculation.", GREEN),
    ("Verifiable Citations", "Exact page numbers, section classifications, and chunk IDs (e.g. doc-CH-0396) for transparent provenance.", NAVY),
    ("Two-Tier Safety Gates", "Regex interceptor for patient-specific dosing requests + score threshold gating for out-of-scope queries.", CRIMSON),
    ("Clinical Speed & Reliability", "FastAPI Linux architecture on Azure App Service with sub-second hybrid retrieval and structured markdown output.", NAVY),
]

for i, (title, desc, color) in enumerate(pillars):
    x = 0.8 + (i % 2) * 5.9
    y = 1.6 + (i // 2) * 2.6
    add_card(s3, x, y, 5.7, 2.3)
    tb = s3.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.25), Inches(5.1), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = color; p.space_after = Pt(6)
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12); p2.font.color.rgb = MUTED

# ══════════════════════════════════════════════════════════════════════
# SLIDE 4: The 3 Indexed Clinical Guidelines
# ══════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
set_slide_background(s4, LIGHT_BG)
add_header(s4, "Knowledge Corpus", "Three Authoritative Guidelines Powering SickleSense")

sources = [
    ("NHLBI (2014)", "National Heart, Lung, and Blood Institute (NIH)", 
     "Comprehensive Evidence-Based Management of SCD",
     "• Hydroxyurea initiation thresholds in adults & children\n• Annual Transcranial Doppler (TCD) stroke screening\n• Acute chest syndrome (ACS) diagnosis & management\n• Transfusion therapy & chronic kidney disease (CKD)"),
    
    ("JCM (2024)", "Journal of Clinical Medicine (MDPI)", 
     "Multicenter Retrospective Analysis Across Age Groups",
     "• Real-world Saudi multicenter cohort data (Taif region)\n• Splenic disease manifestations and acute complications\n• Hospital length-of-stay and ICU admission rates\n• Age-group epidemiological comparison"),
    
    ("WHO (2026)", "World Health Organization (Geneva)", 
     "Consolidated Guidelines for SCD in Children & Adolescents",
     "• Universal newborn screening & diagnostic pathways\n• Penicillin & malaria prophylaxis in endemic regions\n• Universal pediatric hydroxyurea recommendation\n• Growth faltering, danger signs, & referral criteria"),
]

for i, (tag, pub, title, bullets) in enumerate(sources):
    x = 0.8 + i * 4.0
    add_card(s4, x, 1.6, 3.8, 5.2)
    tb = s4.shapes.add_textbox(Inches(x + 0.25), Inches(1.8), Inches(3.3), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = tag
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = CRIMSON; p.space_after = Pt(2)
    p_pub = tf.add_paragraph()
    p_pub.text = pub
    p_pub.font.size = Pt(10); p_pub.font.bold = True; p_pub.font.color.rgb = NAVY; p_pub.space_after = Pt(6)
    p_title = tf.add_paragraph()
    p_title.text = title
    p_title.font.size = Pt(11); p_title.font.italic = True; p_title.font.color.rgb = MUTED; p_title.space_after = Pt(10)
    p_b = tf.add_paragraph()
    p_b.text = bullets
    p_b.font.size = Pt(11); p_b.font.color.rgb = NAVY

# ══════════════════════════════════════════════════════════════════════
# SLIDE 5: Parent-Child Chunking Strategy
# ══════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
set_slide_background(s5, LIGHT_BG)
add_header(s5, "Data Engineering", "Parent-Child Chunking: Precision Without Losing Context")

# Left Box: The Strategy
add_card(s5, 0.8, 1.6, 5.7, 5.2)
tb = s5.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(5.1), Inches(4.8))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Why Parent-Child Chunking?"
p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = NAVY; p.space_after = Pt(10)
p_body = tf.add_paragraph()
p_body.text = (
    "• Small Chunks (800 chars): Optimize dense embedding similarity and pinpoint exact sentences in dense clinical text.\n\n"
    "• Parent References (Page-Level): Preserve chapter context, section hierarchy, and guideline tables.\n\n"
    "• 150-Char Overlap: Guarantees medical terms and multi-word drug protocols (e.g. '15-20 mg/kg/day') are never split across boundaries.\n\n"
    "• Result: High Precision@3 retrieval with complete paragraph context restored during LLM synthesis."
)
p_body.font.size = Pt(12); p_body.font.color.rgb = MUTED

# Right Box: Schema Representation
add_card(s5, 6.8, 1.6, 5.7, 5.2, bg_color=DARK_CARD, border_color=None)
tb = s5.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.1), Inches(4.8))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Indexed Vector Schema (Azure AI Search)"
p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = RGBColor(255, 138, 122); p.space_after = Pt(12)
p_code = tf.add_paragraph()
p_code.text = (
    "{\n"
    "  \"chunk_id\": \"nhlbi-scd-2014-CH-0396\",\n"
    "  \"parent_id\": \"nhlbi-scd-2014-P0077\",\n"
    "  \"document_id\": \"nhlbi-scd-2014\",\n"
    "  \"section\": \"Hydroxyurea Therapy\",\n"
    "  \"page_number\": 77,\n"
    "  \"text\": \"In adults with SCA who have >= 3...\",\n"
    "  \"vector\": [0.0124, -0.0451, ..., 768 dims]\n"
    "}"
)
p_code.font.size = Pt(11); p_code.font.color.rgb = RGBColor(200, 220, 240)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 6: Hybrid Search & Reciprocal Rank Fusion (RRF)
# ══════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
set_slide_background(s6, LIGHT_BG)
add_header(s6, "Retrieval Architecture", "Hybrid Search & Reciprocal Rank Fusion (RRF)")

# 3 Horizontal Step Cards
steps = [
    ("1. Dense Vector Embeddings", "NeuML/biomedbert-base-embeddings", 
     "Transforms clinical queries into 768-dimensional normalized dense vectors. Trained specifically on PubMed and biomedical literature to capture deep clinical semantic relationships."),
    ("2. Lexical BM25 Search", "Exact Medical Acronym Matching", 
     "Performs inverted-index BM25 lexical search to ensure precise matching for critical medical abbreviations (VOC, ACS, TCD, HbSS, HbSβ⁰) that vector search might blur."),
    ("3. Reciprocal Rank Fusion (RRF)", "Server-Side Score Merging", 
     "Azure AI Search merges dense vector and BM25 rankings via RRF algorithm. Scores are calibrated against threshold (0.025) to discard irrelevant questions before generation."),
]

for i, (title, sub, desc) in enumerate(steps):
    x = 0.8 + i * 4.0
    add_card(s6, x, 1.6, 3.8, 5.2)
    tb = s6.shapes.add_textbox(Inches(x + 0.25), Inches(1.8), Inches(3.3), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = NAVY; p.space_after = Pt(3)
    p_sub = tf.add_paragraph()
    p_sub.text = sub
    p_sub.font.size = Pt(10); p_sub.font.bold = True; p_sub.font.color.rgb = CRIMSON; p_sub.space_after = Pt(10)
    p_desc = tf.add_paragraph()
    p_desc.text = desc
    p_desc.font.size = Pt(12); p_desc.font.color.rgb = MUTED

# ══════════════════════════════════════════════════════════════════════
# SLIDE 7: Backend Architecture & Azure Deployment
# ══════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
set_slide_background(s7, LIGHT_BG)
add_header(s7, "Cloud & Backend Infrastructure", "FastAPI on Azure App Service & Azure AI Search")

# Left Column: Backend Engine
add_card(s7, 0.8, 1.6, 5.7, 5.2)
tb = s7.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(5.1), Inches(4.8))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "FastAPI Backend Architecture"
p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = NAVY; p.space_after = Pt(10)
p_body = tf.add_paragraph()
p_body.text = (
    "• High-Performance Async API: Decoupled Python service built with FastAPI and Pydantic validation.\n\n"
    "• Dedicated App Service Linux: Configured with Always-On to eliminate cold starts typical of serverless functions.\n\n"
    "• Groq LLM Acceleration: Queries openai/gpt-oss-120b in sub-second response times under strict temperature=0 grounding.\n\n"
    "• Automated OpenAPI Docs: Live interactive Swagger UI at /docs and monitoring at /health."
)
p_body.font.size = Pt(12); p_body.font.color.rgb = MUTED

# Right Column: Azure Services
add_card(s7, 6.8, 1.6, 5.7, 5.2)
tb = s7.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.1), Inches(4.8))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Azure Cloud Infrastructure"
p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = NAVY; p.space_after = Pt(10)
p_body = tf.add_paragraph()
p_body.text = (
    "• Azure AI Search Index (creativa-hackathon-pc): Hosts 768d vector index with HNSW cosine similarity alongside BM25 indexes.\n\n"
    "• Azure App Service (sicklesense-api): Production Linux container hosting the FastAPI backend in Sweden Central region.\n\n"
    "• CORS & Security: Configured for secure frontend access without exposing database credentials or API keys.\n\n"
    "• Zero-Downtime Reliability: Fully managed scalable cloud instance."
)
p_body.font.size = Pt(12); p_body.font.color.rgb = MUTED

# ══════════════════════════════════════════════════════════════════════
# SLIDE 8: Frontend Integration & Deployment Pipeline
# ══════════════════════════════════════════════
s8 = prs.slides.add_slide(blank_layout)
set_slide_background(s8, LIGHT_BG)
add_header(s8, "Frontend & CI/CD", "Vite + React SPA & GitHub Actions Deployment")

# Grid: 3 Cards
cards = [
    ("Modern React UI", "Vite + React + Lucide", 
     "• Navy + Crimson clinical branding matching SickleSense mascot.\n• Markdown tables with remark-gfm responsive containers.\n• Live backend health checker & interactive FAQ accordion.\n• Zero chat clutter: structured evidence notes."),
    ("Decoupled Integration", "RESTful JSON Protocol", 
     "• POST /api/query: Sends { query, top_k } payload.\n• Dynamic citation mapping displaying score, section, & chunk ID.\n• AbortController timeout handling (30s resilience).\n• Direct sample query loading from Sources page."),
    ("Automated CI/CD", "GitHub Actions & Pages", 
     "• Automated deployment workflow on git push to master.\n• Production Vite build outputting optimized static assets.\n• Relative base path configuration for reliable subpath routing.\n• Zero hosting maintenance overhead."),
]

for i, (title, sub, desc) in enumerate(cards):
    x = 0.8 + i * 4.0
    add_card(s8, x, 1.6, 3.8, 5.2)
    tb = s8.shapes.add_textbox(Inches(x + 0.25), Inches(1.8), Inches(3.3), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = NAVY; p.space_after = Pt(3)
    p_sub = tf.add_paragraph()
    p_sub.text = sub
    p_sub.font.size = Pt(10); p_sub.font.bold = True; p_sub.font.color.rgb = CRIMSON; p_sub.space_after = Pt(10)
    p_desc = tf.add_paragraph()
    p_desc.text = desc
    p_desc.font.size = Pt(11); p_desc.font.color.rgb = MUTED

# ══════════════════════════════════════════════════════════════════════
# SLIDE 9: Evaluation & Measured Benchmarks
# ══════════════════════════════════════════════
s9 = prs.slides.add_slide(blank_layout)
set_slide_background(s9, LIGHT_BG)
add_header(s9, "Performance & Verification", "Empirical Evaluation Across 20 Clinical Benchmark Prompts")

# Stat boxes
stat_boxes = [
    ("Precision@3", "0.500", "50% of top-3 retrieved passages directly answer clinical intent", GREEN),
    ("Precision@5", "0.463", "Broad clinical context retained across 5 retrieved chunks", NAVY),
    ("Out-of-Scope Rejection", "100%", "0.000 false-relevance rate on unrelated medical questions", CRIMSON),
]

for i, (label, val, note, color) in enumerate(stat_boxes):
    x = 0.8 + i * 4.0
    add_card(s9, x, 1.6, 3.8, 2.3)
    tb = s9.shapes.add_textbox(Inches(x + 0.2), Inches(1.75), Inches(3.4), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = MUTED
    p_val = tf.add_paragraph()
    p_val.text = val
    p_val.font.size = Pt(36); p_val.font.bold = True; p_val.font.color.rgb = color
    p_note = tf.add_paragraph()
    p_note.text = note
    p_note.font.size = Pt(10); p_note.font.color.rgb = MUTED

# Bottom Box: Benchmark breakdown
add_card(s9, 0.8, 4.2, 11.733, 2.6)
tb = s9.shapes.add_textbox(Inches(1.1), Inches(4.35), Inches(11.1), Inches(2.3))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Benchmark Categories Tested During Development"
p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = NAVY; p.space_after = Pt(6)
p_body = tf.add_paragraph()
p_body.text = (
    "1. Direct Guideline Fact Lookup (e.g. 'When should hydroxyurea be started in adults?') -> Exact guideline citation matched.\n"
    "2. Paraphrased Clinical Queries (e.g. 'How to lower alloimmunization risk from transfusion?') -> Correct matching across chapters.\n"
    "3. Thresholds & Acronyms (e.g. 'What does TCD stand for and how is stroke screened?') -> Precise abbreviation expansion.\n"
    "4. Out-of-Scope Disease Questions (e.g. 'Ibuprofen dosage for migraine in healthy adult') -> Gated & safely refused.\n"
    "5. Unsafe Patient-Specific Dosage Inquiries -> Intercepted immediately before LLM invocation."
)
p_body.font.size = Pt(11); p_body.font.color.rgb = MUTED

# ══════════════════════════════════════════════════════════════════════
# SLIDE 10: Conclusion & Team Attribution (Dark Navy)
# ══════════════════════════════════════════════
s10 = prs.slides.add_slide(blank_layout)
set_slide_background(s10, NAVY)

tb10 = s10.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
tf10 = tb10.text_frame
tf10.word_wrap = True

p = tf10.paragraphs[0]
p.text = "SICKLESENSE"
p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = WHITE; p.space_after = Pt(12)

p_tag = tf10.add_paragraph()
p_tag.text = "\"Sources Before Certainty\""
p_tag.font.size = Pt(22); p_tag.font.italic = True; p_tag.font.color.rgb = RGBColor(255, 138, 122); p_tag.space_after = Pt(24)

p_sum = tf10.add_paragraph()
p_sum.text = (
    "SickleSense demonstrates how domain-adapted embeddings, parent-child chunking, and hybrid RRF search "
    "can transform hundreds of pages of complex medical guidelines into an instant, traceable, and safe clinical assistant.\n\n"
    "Live Web App: https://O-076.github.io/sicklesense/\n"
    "FastAPI Swagger Docs: https://sicklesense-api-dgh3aqdwb3eghdc4.swedencentral-01.azurewebsites.net/docs"
)
p_sum.font.size = Pt(13); p_sum.font.color.rgb = RGBColor(203, 214, 226); p_sum.space_after = Pt(28)

p_team = tf10.add_paragraph()
p_team.text = "Developed with care by Team MINOR THREAT 2026 for the Creativa SCD AI Hackathon"
p_team.font.size = Pt(12); p_team.font.bold = True; p_team.font.color.rgb = RGBColor(160, 180, 205)

output_path = "SickleSense_Pitch_Deck.pptx"
prs.save(output_path)
print(f"Presentation saved successfully to {output_path}")
